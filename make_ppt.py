"""ECG RAG 시스템 발표자료 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── 컬러 팔레트 ────────────────────────────────────────────────────────────────
C_BG       = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT   = RGBColor(0x00, 0xB4, 0xD8)
C_ACCENT2  = RGBColor(0x90, 0xE0, 0xEF)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY     = RGBColor(0xAA, 0xBB, 0xCC)
C_CARD     = RGBColor(0x1A, 0x2E, 0x42)
C_GREEN    = RGBColor(0x2D, 0xD4, 0x8E)
C_ORANGE   = RGBColor(0xF7, 0x8C, 0x6C)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ──────────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color=C_BG):
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def txt(slide, text, x, y, w, h, size=24, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color

def accent_bar(slide, y=Inches(0.55)):
    bar = slide.shapes.add_shape(1, Inches(0.55), y, Inches(2.4), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()

def card(slide, x, y, w, h, color=C_CARD, border=C_ACCENT):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)

def section_tag(slide, label):
    tag = slide.shapes.add_shape(1, Inches(0.55), Inches(6.8), Inches(2.4), Inches(0.45))
    tag.fill.solid()
    tag.fill.fore_color.rgb = C_ACCENT
    tag.line.fill.background()
    t = tag.text_frame.paragraphs[0]
    t.alignment = PP_ALIGN.CENTER
    r = t.add_run()
    r.text = label
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = C_BG

def bullets(slide, items, x, y, w, h, size=16, color=C_WHITE):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# 01 — 타이틀
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)

deco = sl.shapes.add_shape(1, Inches(8.8), Inches(-0.5), Inches(5.0), Inches(8.5))
deco.fill.solid()
deco.fill.fore_color.rgb = RGBColor(0x10, 0x26, 0x3B)
deco.line.fill.background()
line = sl.shapes.add_shape(1, Inches(8.77), Inches(0), Pt(3), H)
line.fill.solid()
line.fill.fore_color.rgb = C_ACCENT
line.line.fill.background()

txt(sl, "ECG RAG System", Inches(0.7), Inches(1.6), Inches(7.8), Inches(1.2),
    size=52, bold=True)
txt(sl, "심전도 AI 진단 설명 시스템", Inches(0.7), Inches(2.9), Inches(7.8), Inches(0.7),
    size=28, color=C_ACCENT2)
txt(sl, "RAG 기반 의학 지식 검색 · 자동 설명 생성", Inches(0.7), Inches(3.7), Inches(7.8), Inches(0.6),
    size=20, color=C_GRAY)
txt(sl, "2026.05", Inches(9.2), Inches(3.5), Inches(3.5), Inches(0.5),
    size=18, color=C_GRAY, align=PP_ALIGN.CENTER)

for i, kw in enumerate(["PDF 교재 7권", "6,055 청크", "BGE-M3 임베딩", "Qdrant 벡터 DB"]):
    y = Inches(1.8 + i * 0.9)
    dot = sl.shapes.add_shape(9, Inches(9.1), y + Inches(0.12), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = C_ACCENT
    dot.line.fill.background()
    txt(sl, kw, Inches(9.45), y, Inches(3.3), Inches(0.55), size=16)

# ══════════════════════════════════════════════════════════════════════════════
# 02 — 왜 RAG인가?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "왜 RAG인가?", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7), size=34, bold=True)
section_tag(sl, "BACKGROUND")

card(sl, Inches(0.55), Inches(1.6), Inches(5.8), Inches(4.8))
txt(sl, "😟  기존 방식의 한계", Inches(0.8), Inches(1.75), Inches(5.3), Inches(0.55),
    size=17, bold=True, color=C_ACCENT)
bullets(sl, [
    "• 심전도 분석기는 진단명만 출력",
    "• 의사/환자에게 근거 설명이 없음",
    "• 교재 지식을 코드에 하드코딩하면\n  유지보수 불가능",
    "• LLM만 단독으로 쓰면 의학적\n  환각(Hallucination) 위험",
], Inches(0.8), Inches(2.4), Inches(5.3), Inches(3.5))

card(sl, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.8))
txt(sl, "✅  RAG로 해결", Inches(7.05), Inches(1.75), Inches(5.3), Inches(0.55),
    size=17, bold=True, color=C_GREEN)
bullets(sl, [
    "• 실제 의학 교재를 검색 소스로 활용",
    "• XML 진단 결과 → 관련 챕터 자동 검색",
    "• LLM이 검색된 근거로 설명 생성",
    "• 교재 추가만으로 지식 확장 가능",
    "• 출처(책/페이지) 제공으로 신뢰도 확보",
], Inches(7.05), Inches(2.4), Inches(5.3), Inches(3.5))

# ══════════════════════════════════════════════════════════════════════════════
# 03 — RAG 개념
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "RAG란 무엇인가?", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7), size=34, bold=True)
section_tag(sl, "CONCEPT")
txt(sl, "Retrieval-Augmented Generation", Inches(0.55), Inches(1.55), Inches(12), Inches(0.55),
    size=22, color=C_ACCENT2, bold=True)

labels = ["1. 검색 (Retrieve)", "2. 증강 (Augment)", "3. 생성 (Generate)"]
descs  = [
    "질문과 관련된\n문서를 벡터 DB에서\n유사도 검색",
    "검색된 문서를\n프롬프트에 추가해\nLLM에게 컨텍스트 제공",
    "LLM이 검색된\n근거를 바탕으로\n신뢰할 수 있는 답변 생성",
]
bg_cols = [RGBColor(0x0A,0x3D,0x5C), RGBColor(0x0A,0x4A,0x3A), RGBColor(0x3A,0x2A,0x5C)]

for i, (label, desc, bg) in enumerate(zip(labels, descs, bg_cols)):
    x = Inches(0.55 + i * 4.2)
    c = sl.shapes.add_shape(1, x, Inches(2.5), Inches(3.9), Inches(3.2))
    c.fill.solid(); c.fill.fore_color.rgb = bg
    c.line.color.rgb = C_ACCENT; c.line.width = Pt(1)
    txt(sl, label, x + Inches(0.2), Inches(2.65), Inches(3.5), Inches(0.55),
        size=18, bold=True, color=C_ACCENT)
    txt(sl, desc, x + Inches(0.2), Inches(3.3), Inches(3.5), Inches(2.0), size=16)
    if i < 2:
        txt(sl, "→", Inches(4.3 + i * 4.2), Inches(3.8), Inches(0.5), Inches(0.5),
            size=28, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

txt(sl, "핵심: LLM의 창작이 아닌, 실제 문서에서 찾아서 설명하는 방식",
    Inches(0.55), Inches(6.1), Inches(12), Inches(0.6),
    size=19, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 04 — 시스템 전체 구조
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "시스템 전체 구조", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7), size=34, bold=True)
section_tag(sl, "ARCHITECTURE")

txt(sl, "① 사전 처리 (Ingest Pipeline)  —  1회 실행",
    Inches(0.55), Inches(1.6), Inches(12), Inches(0.5), size=15, bold=True, color=C_ACCENT)

for i, s in enumerate(["PDF 교재\n(7권)", "텍스트\n추출", "의미 단위\n청킹", "BGE-M3\n임베딩", "Qdrant\n저장"]):
    x = Inches(0.55 + i * 2.45)
    c = sl.shapes.add_shape(1, x, Inches(2.15), Inches(2.1), Inches(0.85))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x0A,0x3A,0x5A)
    c.line.color.rgb = C_ACCENT; c.line.width = Pt(0.75)
    txt(sl, s, x, Inches(2.15), Inches(2.1), Inches(0.85),
        size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(sl, "▶", Inches(2.5 + i * 2.45), Inches(2.3), Inches(0.4), Inches(0.5),
            size=14, color=C_ACCENT, align=PP_ALIGN.CENTER)

div = sl.shapes.add_shape(1, Inches(0.4), Inches(3.3), Inches(12.5), Pt(1))
div.fill.solid(); div.fill.fore_color.rgb = C_GRAY; div.line.fill.background()

txt(sl, "② 실시간 쿼리 Pipeline  —  매 요청마다 실행",
    Inches(0.55), Inches(3.4), Inches(12), Inches(0.5), size=15, bold=True, color=C_GREEN)

for i, s in enumerate(["ECG XML\n분석 결과", "XML\n파싱", "쿼리\n임베딩", "Qdrant\n검색", "LLM\n설명 생성", "진단\n설명 출력"]):
    x = Inches(0.55 + i * 2.1)
    c = sl.shapes.add_shape(1, x, Inches(3.9), Inches(1.9), Inches(0.85))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x0A,0x3D,0x2A)
    c.line.color.rgb = C_GREEN; c.line.width = Pt(0.75)
    txt(sl, s, x, Inches(3.9), Inches(1.9), Inches(0.85),
        size=13, color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        txt(sl, "▶", Inches(2.3 + i * 2.1), Inches(4.03), Inches(0.35), Inches(0.5),
            size=13, color=C_GREEN, align=PP_ALIGN.CENTER)

txt(sl, "↕  공유 벡터 DB", Inches(6.0), Inches(3.1), Inches(2.5), Inches(0.4),
    size=12, color=C_ACCENT2, align=PP_ALIGN.CENTER)

txt(sl, "인제스트는 1회 실행, 쿼리는 매 요청마다 실행  |  두 파이프라인이 Qdrant를 공유",
    Inches(0.55), Inches(5.2), Inches(12.2), Inches(0.5),
    size=16, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 05 — 데이터 소스
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "데이터 소스: 심전도 교재 7권", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "DATA SOURCE")

books = [
    ("만화로보는 심전도",           "한국어", "기초",    "97"),
    ("Basic Concepts of EKG",     "영어",   "기초",   "405"),
    ("Goldberger's Clinical ECG", "영어",   "임상",  "1,503"),
    ("Marriott's 13판 (번역)",     "한국어", "임상",  "2,384"),
    ("Marriott's 13판 (원문)",     "영어",   "임상",   "946"),
    ("Basic Electrocardiography", "영어",   "임상",   "225"),
    ("Inherited Arrhythmias",     "영어",   "전문",   "495"),
]
hcols   = [Inches(0.55), Inches(5.4), Inches(7.2), Inches(9.0)]
hwidths = [Inches(4.7),  Inches(1.6), Inches(1.6), Inches(1.6)]
for hx, hw, hl in zip(hcols, hwidths, ["교재명", "언어", "레벨", "청크 수"]):
    hb = sl.shapes.add_shape(1, hx, Inches(1.55), hw, Inches(0.45))
    hb.fill.solid(); hb.fill.fore_color.rgb = C_ACCENT; hb.line.fill.background()
    txt(sl, hl, hx + Inches(0.1), Inches(1.6), hw, Inches(0.4),
        size=14, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

lc_map = {"기초": C_GREEN, "임상": C_ACCENT, "전문": C_ORANGE}
for i, (name, lang, level, cnt) in enumerate(books):
    y = Inches(2.1 + i * 0.62)
    row = sl.shapes.add_shape(1, Inches(0.55), y, Inches(12.2), Inches(0.57))
    row.fill.solid()
    row.fill.fore_color.rgb = RGBColor(0x13,0x24,0x38) if i%2==0 else RGBColor(0x17,0x2C,0x44)
    row.line.fill.background()
    txt(sl, name, Inches(0.7),  y+Inches(0.08), Inches(4.5), Inches(0.45), size=15)
    txt(sl, lang, Inches(5.4),  y+Inches(0.08), Inches(1.5), Inches(0.45), size=15, color=C_GRAY, align=PP_ALIGN.CENTER)
    lb = sl.shapes.add_shape(1, Inches(7.1), y+Inches(0.1), Inches(1.2), Inches(0.35))
    lb.fill.solid(); lb.fill.fore_color.rgb = lc_map.get(level, C_WHITE); lb.line.fill.background()
    lt = lb.text_frame.paragraphs[0]; lt.alignment = PP_ALIGN.CENTER
    lr = lt.add_run(); lr.text = level; lr.font.size = Pt(12); lr.font.bold = True; lr.font.color.rgb = C_BG
    txt(sl, cnt, Inches(9.0), y+Inches(0.08), Inches(1.5), Inches(0.45),
        size=15, color=C_ACCENT2, align=PP_ALIGN.CENTER, bold=True)

txt(sl, "총 6,055 청크  |  한국어 2,481 · 영어 3,574",
    Inches(0.55), Inches(6.6), Inches(12.2), Inches(0.5),
    size=17, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 06 — PDF 추출
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "Step 1 — PDF 텍스트 추출", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "INGEST · STEP 1")
txt(sl, "도구: PyMuPDF  |  208MB 교재 → 페이지별 순수 텍스트",
    Inches(0.55), Inches(1.55), Inches(12), Inches(0.5), size=18, color=C_ACCENT2)

for i, (title, desc) in enumerate([
    ("블록 단위 추출", "PDF의 텍스트 블록 구조를\n그대로 활용해 단락 경계 보존"),
    ("헤더/푸터 제거", "페이지 상하 7% 영역을\n자동으로 필터링"),
    ("짧은 페이지 스킵", "80자 미만 페이지 제외\n(그림 전용, 빈 페이지)"),
]):
    x = Inches(0.55 + i * 4.2)
    card(sl, x, Inches(2.4), Inches(3.9), Inches(2.5))
    txt(sl, title, x+Inches(0.2), Inches(2.6), Inches(3.5), Inches(0.55),
        size=18, bold=True, color=C_ACCENT)
    txt(sl, desc, x+Inches(0.2), Inches(3.2), Inches(3.5), Inches(1.5), size=16)

card(sl, Inches(0.55), Inches(5.1), Inches(12.2), Inches(1.2))
txt(sl, "출력 구조 (페이지 1개)",
    Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.45), size=15, bold=True, color=C_ACCENT)
txt(sl, '{ "page_num": 12,  "text": "Sinus rhythm requires ...",  "source_file": "basic-concepts-of-ekg.pdf" }',
    Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.5), size=14, color=C_ACCENT2, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# 07 — 청킹
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "Step 2 — 의미 경계 청킹", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "INGEST · STEP 2")
txt(sl, "고정 크기 자르기(✗)  →  의학 문맥을 이해한 경계 분리(✓)",
    Inches(0.55), Inches(1.55), Inches(12), Inches(0.5), size=18, color=C_ACCENT2)

card(sl, Inches(0.55), Inches(2.2), Inches(5.8), Inches(4.4))
txt(sl, "청킹 규칙", Inches(0.8), Inches(2.35), Inches(5.2), Inches(0.5),
    size=18, bold=True, color=C_ACCENT)
bullets(sl, [
    "① 섹션 제목 감지 → 새 청크 시작",
    "② 진단 기준 블록(ms/mm/bpm 포함)",
    "   → 절대 중간에 자르지 않음",
    "③ 목록(bullet/번호) → 통째로 유지",
    "④ 목표 500 토큰, 최대 700 토큰",
    "⑤ 2문장 오버랩으로 문맥 연속성 유지",
], Inches(0.8), Inches(2.95), Inches(5.2), Inches(3.3))

card(sl, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.4))
txt(sl, "청크 타입 분류", Inches(7.05), Inches(2.35), Inches(5.2), Inches(0.5),
    size=18, bold=True, color=C_ACCENT)
for i, (tag, desc, col) in enumerate([
    ("criteria",  "진단 기준, 수치 임계값 포함\n→ LLM에게 정밀 근거 제공", C_GREEN),
    ("narrative", "임상 설명, 배경 지식 텍스트\n→ LLM에게 문맥 제공",     C_ACCENT),
]):
    ty = Inches(3.0 + i * 1.7)
    tb = sl.shapes.add_shape(1, Inches(7.05), ty, Inches(1.3), Inches(0.45))
    tb.fill.solid(); tb.fill.fore_color.rgb = col; tb.line.fill.background()
    tt = tb.text_frame.paragraphs[0]; tt.alignment = PP_ALIGN.CENTER
    tr = tt.add_run(); tr.text = tag; tr.font.size = Pt(13); tr.font.bold = True; tr.font.color.rgb = C_BG
    txt(sl, desc, Inches(8.5), ty, Inches(3.8), Inches(0.85), size=15)

txt(sl, "결과: criteria 1,931개  +  narrative 4,124개  =  총 6,055 청크",
    Inches(6.8), Inches(5.8), Inches(5.8), Inches(0.55),
    size=15, bold=True, color=C_ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# 08 — 임베딩
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "Step 3 — BGE-M3 임베딩", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "INGEST · STEP 3")
txt(sl, "텍스트 → 숫자 벡터 변환  |  의미가 비슷한 텍스트는 가까운 벡터",
    Inches(0.55), Inches(1.55), Inches(12), Inches(0.5), size=18, color=C_ACCENT2)

txt(sl, "임베딩이란?", Inches(0.55), Inches(2.2), Inches(12), Inches(0.5),
    size=20, bold=True)
for i, (src, vec, note) in enumerate([
    ('"Atrial Fibrillation"',    "[0.23, -0.87, 0.41 ...]", ""),
    ('"심방세동 진단 기준"',         "[0.21, -0.84, 0.39 ...]", "← 거의 같은 방향!"),
    ('"심전도 P파 부재"',            "[0.19, -0.80, 0.44 ...]", "← 유사"),
]):
    y = Inches(2.85 + i * 0.75)
    card(sl, Inches(0.55), y, Inches(11.8), Inches(0.65))
    txt(sl, src,  Inches(0.75), y+Inches(0.1), Inches(3.5), Inches(0.5), size=16, color=C_ACCENT2, italic=True)
    txt(sl, "→",  Inches(4.3),  y+Inches(0.1), Inches(0.5), Inches(0.5), size=18, color=C_GRAY)
    txt(sl, vec,  Inches(4.8),  y+Inches(0.1), Inches(4.5), Inches(0.5), size=14, color=C_GREEN, italic=True)
    if note:
        txt(sl, note, Inches(9.4), y+Inches(0.1), Inches(2.7), Inches(0.5), size=13, color=C_ACCENT, italic=True)

txt(sl, "BGE-M3 선택 이유", Inches(0.55), Inches(5.2), Inches(12), Inches(0.45),
    size=18, bold=True)
for i, (title, desc) in enumerate([
    ("한/영 동시 지원",     "한국어 번역본·영어 원문을 같은 벡터 공간에서 검색"),
    ("1,024차원 Dense",    "높은 표현력, 의학 용어 미묘한 차이 구분"),
    ("오픈소스·자체 호스팅", "외부 API 의존 없음, 비용 없음"),
]):
    x = Inches(0.55 + i * 4.2)
    card(sl, x, Inches(5.7), Inches(3.9), Inches(1.3))
    txt(sl, title, x+Inches(0.15), Inches(5.75), Inches(3.6), Inches(0.45), size=15, bold=True, color=C_ACCENT)
    txt(sl, desc,  x+Inches(0.15), Inches(6.2),  Inches(3.6), Inches(0.7),  size=13, color=C_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# 09 — Qdrant
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "Step 4 — Qdrant 벡터 DB", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "INGEST · STEP 4")

card(sl, Inches(0.55), Inches(1.6), Inches(5.8), Inches(5.3))
txt(sl, "Qdrant란?", Inches(0.8), Inches(1.75), Inches(5.3), Inches(0.5), size=18, bold=True, color=C_ACCENT)
bullets(sl, [
    "• 벡터 유사도 검색에 특화된 DB",
    "• 고속 코사인 유사도 계산",
    "• 메타데이터 필터링 동시 지원",
    "• 로컬 파일 모드 → 서버 불필요",
    "• 추후 서버 모드로 그대로 전환 가능",
], Inches(0.8), Inches(2.3), Inches(5.3), Inches(3.8))

card(sl, Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.3))
txt(sl, "저장 구조 (청크 1개)", Inches(7.05), Inches(1.75), Inches(5.3), Inches(0.5),
    size=18, bold=True, color=C_ACCENT)
for i, (field, val) in enumerate([
    ("id",             "고유 UUID"),
    ("vector",         "1024차원 숫자 배열"),
    ("content",        "청크 원문 텍스트"),
    ("source_book",    "marriotts_13ed_ko 등"),
    ("language",       "en / ko"),
    ("audience_level", "basic / clinical / specialist"),
    ("content_type",   "criteria / narrative"),
    ("page_num",       "원본 페이지 번호"),
]):
    y = Inches(2.35 + i * 0.52)
    row = sl.shapes.add_shape(1, Inches(6.8), y, Inches(5.8), Inches(0.48))
    row.fill.solid()
    row.fill.fore_color.rgb = RGBColor(0x13,0x26,0x3E) if i%2==0 else C_BG
    row.line.fill.background()
    txt(sl, field, Inches(7.0), y+Inches(0.05), Inches(2.0), Inches(0.4),
        size=13, color=C_ACCENT2, bold=True, italic=True)
    txt(sl, val,   Inches(9.1), y+Inches(0.05), Inches(3.2), Inches(0.4), size=13, color=C_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# 10 — 메타데이터 필터링
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "메타데이터 필터링", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "METADATA")
txt(sl, "벡터 유사도 + 조건 필터 = 정밀한 검색",
    Inches(0.55), Inches(1.55), Inches(12), Inches(0.5), size=20, color=C_ACCENT2, bold=True)

for i, (title, filt, result) in enumerate([
    ("환자 대상 설명", "audience_level = basic\nlanguage = ko",
     "만화 심전도 · 기초 교재에서\n쉬운 설명 검색"),
    ("의사 대상 설명", "audience_level = clinical\nlanguage = en/ko",
     "Marriott's · Goldberger's에서\n임상 기준 검색"),
    ("전문의 대상", "audience_level = specialist",
     "Inherited Arrhythmias에서\n유전성 부정맥 전문 내용 검색"),
]):
    x = Inches(0.55 + i * 4.2)
    card(sl, x, Inches(2.3), Inches(3.9), Inches(3.8))
    txt(sl, title, x+Inches(0.2), Inches(2.45), Inches(3.5), Inches(0.5), size=17, bold=True, color=C_ACCENT)
    txt(sl, "필터:", x+Inches(0.2), Inches(3.0),  Inches(3.5), Inches(0.4), size=13, color=C_GRAY)
    txt(sl, filt,   x+Inches(0.2), Inches(3.4),  Inches(3.5), Inches(0.7), size=14, color=C_GREEN, italic=True)
    txt(sl, "결과:", x+Inches(0.2), Inches(4.15), Inches(3.5), Inches(0.4), size=13, color=C_GRAY)
    txt(sl, result, x+Inches(0.2), Inches(4.55), Inches(3.5), Inches(0.9), size=14)

txt(sl, "같은 쿼리라도 대상에 따라 다른 교재에서 다른 수준의 설명을 검색",
    Inches(0.55), Inches(6.3), Inches(12.2), Inches(0.5),
    size=17, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 11 — XML 파싱
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "쿼리 Step 1 — ECG XML 파싱", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "QUERY · STEP 1")
txt(sl, "심전도 분석기 출력(XML) → RAG 쿼리 컨텍스트 자동 변환",
    Inches(0.55), Inches(1.55), Inches(12), Inches(0.5), size=18, color=C_ACCENT2)

card(sl, Inches(0.55), Inches(2.1), Inches(5.8), Inches(4.6))
txt(sl, "ECG 분석기 XML 출력", Inches(0.8), Inches(2.25), Inches(5.2), Inches(0.45),
    size=16, bold=True, color=C_ACCENT)
bullets(sl, [
    "<ECGReport>",
    "  <Measurements>",
    "    <HeartRate unit='bpm'>148</HeartRate>",
    "    <PRInterval unit='ms'>absent</PRInterval>",
    "    <QRSDuration unit='ms'>88</QRSDuration>",
    "  </Measurements>",
    "  <Diagnoses>",
    "    <Diagnosis>Atrial Fibrillation</Diagnosis>",
    "  </Diagnoses>",
    "</ECGReport>",
], Inches(0.8), Inches(2.75), Inches(5.2), Inches(3.5), size=13, color=C_ACCENT2)

card(sl, Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.6))
txt(sl, "파싱 결과 (ECGQueryContext)", Inches(7.05), Inches(2.25), Inches(5.2), Inches(0.45),
    size=16, bold=True, color=C_GREEN)
bullets(sl, [
    "diagnoses:",
    "  ['Atrial Fibrillation']",
    "",
    "abnormal_params:",
    "  ['HeartRate 148bpm (high)']",
    "",
    "natural_query:",
    "  'Atrial Fibrillation HeartRate 148bpm",
    "   ECG diagnosis criteria...'",
], Inches(7.05), Inches(2.75), Inches(5.2), Inches(3.5), size=13)

txt(sl, "정상 범위(HR 60-100bpm, PR 120-200ms 등) 자동 비교 → 이상 파라미터 자동 강조",
    Inches(0.55), Inches(6.85), Inches(12.2), Inches(0.45),
    size=15, color=C_ACCENT2, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 12 — 벡터 검색
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "쿼리 Step 2 — 벡터 검색", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "QUERY · STEP 2")

for i, label in enumerate(["자연어 쿼리\n생성", "BGE-M3\n임베딩", "Qdrant\n코사인 검색", "Top-K\n청크 반환"]):
    x = Inches(0.6 + i * 3.1)
    c = sl.shapes.add_shape(1, x, Inches(1.7), Inches(2.6), Inches(1.1))
    c.fill.solid(); c.fill.fore_color.rgb = C_CARD
    c.line.color.rgb = C_ACCENT; c.line.width = Pt(1.5)
    txt(sl, label, x, Inches(1.7), Inches(2.6), Inches(1.1),
        size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    if i < 3:
        txt(sl, "→", Inches(3.05 + i * 3.1), Inches(2.0), Inches(0.6), Inches(0.6),
            size=26, bold=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

txt(sl, "코사인 유사도란?", Inches(0.55), Inches(3.2), Inches(12), Inches(0.5),
    size=18, bold=True)
txt(sl, "두 벡터가 얼마나 같은 방향을 가리키는지를 -1 ~ 1 사이 점수로 측정  (1에 가까울수록 의미가 유사)",
    Inches(0.55), Inches(3.75), Inches(12.2), Inches(0.5), size=16, color=C_GRAY)

txt(sl, "검색 결과 예시", Inches(0.55), Inches(4.4), Inches(12), Inches(0.45),
    size=18, bold=True)
for i, (score, source, content) in enumerate([
    ("0.712", "marriotts_13ed_ko  p.342  [criteria]",
     "심방세동의 ECG 기준: 불규칙한 RR 간격, P파 소실, 세동파(f파)..."),
    ("0.684", "goldbergers_10ed   p.218  [criteria]",
     "Atrial fibrillation: absent P waves, irregularly irregular rhythm..."),
    ("0.651", "marriotts_13ed_en  p.330  [narrative]",
     "AF is the most common sustained cardiac arrhythmia, occurring..."),
]):
    y = Inches(4.95 + i * 0.62)
    row = sl.shapes.add_shape(1, Inches(0.55), y, Inches(12.2), Inches(0.57))
    row.fill.solid()
    row.fill.fore_color.rgb = RGBColor(0x12,0x24,0x38) if i%2==0 else C_BG
    row.line.fill.background()
    txt(sl, score,   Inches(0.7),  y+Inches(0.08), Inches(0.8), Inches(0.42), size=14, color=C_GREEN, bold=True)
    txt(sl, source,  Inches(1.6),  y+Inches(0.08), Inches(3.5), Inches(0.42), size=12, color=C_ACCENT2)
    txt(sl, content, Inches(5.2),  y+Inches(0.08), Inches(7.3), Inches(0.42), size=12, color=C_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# 13 — LLM 설명 생성
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "쿼리 Step 3 — LLM 설명 생성", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "QUERY · STEP 3")

card(sl, Inches(0.55), Inches(1.6), Inches(12.2), Inches(2.8))
txt(sl, "LLM에게 전달되는 프롬프트 구조", Inches(0.8), Inches(1.75), Inches(11.5), Inches(0.45),
    size=16, bold=True, color=C_ACCENT)
bullets(sl, [
    "[시스템]  당신은 심전도 전문의입니다. 아래 교재 내용을 근거로 설명하세요.",
    "",
    "[교재 검색 결과]",
    "  Source 1: marriotts_13ed_ko p.342 — 심방세동의 ECG 기준: 불규칙한 RR 간격...",
    "  Source 2: goldbergers_10ed  p.218 — AF: absent P waves, irregularly irregular...",
    "",
    "[ECG 분석 결과]  진단: Atrial Fibrillation  |  이상: HR 148bpm (high)",
    "",
    "[질문]  위 심전도 소견에 대해 임상의 수준으로 설명해 주세요.",
], Inches(0.8), Inches(2.25), Inches(11.5), Inches(2.0), size=13)

txt(sl, "↓  LLM 응답", Inches(5.8), Inches(4.55), Inches(2.5), Inches(0.45),
    size=16, color=C_ACCENT2, bold=True, align=PP_ALIGN.CENTER)

card(sl, Inches(0.55), Inches(5.05), Inches(12.2), Inches(1.6))
txt(sl, "심방세동(Atrial Fibrillation)은 심방의 불규칙한 전기 활동으로 발생합니다. ECG상 P파가 소실되고\n"
        "불규칙한 RR 간격이 특징입니다. 현재 심박수 148bpm으로 빠른 심실 반응을 동반하고 있어 즉각적인\n"
        "속도 조절 치료가 필요합니다.  (출처: Marriott's 13판 p.342, Goldberger's 10판 p.218)",
    Inches(0.8), Inches(5.15), Inches(11.5), Inches(1.4), size=14)

# ══════════════════════════════════════════════════════════════════════════════
# 14 — 발견한 문제 & 해결
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "실제 운영 중 발견한 문제 & 해결", Inches(0.55), Inches(0.7), Inches(11), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "LESSONS LEARNED")

for i, (title, prob, sol, col) in enumerate([
    ("① 한국어 파일명 문제",
     "NFD(자모 분해형) vs NFC(조합형) 유니코드\n불일치로 한국어 PDF 2권 인제스트 스킵",
     "unicodedata.normalize('NFC')로\n파일명 정규화 후 해결",
     C_ACCENT),
    ("② 책 색인 페이지 혼입",
     "책 뒷부분 색인이 ECG 용어 밀집\n→ 관련도 높은 것처럼 검색됨",
     "색인 페이지 패턴 감지 필터 추가\n(Phase 2 예정)",
     C_ORANGE),
    ("③ qdrant-client API 변경",
     "v1.14+ 에서 client.search() 제거\n→ AttributeError 발생",
     "client.query_points()로\n마이그레이션 완료",
     C_GREEN),
]):
    y = Inches(1.7 + i * 1.65)
    card(sl, Inches(0.55), y, Inches(12.2), Inches(1.5))
    dot = sl.shapes.add_shape(9, Inches(0.75), y+Inches(0.55), Inches(0.25), Inches(0.25))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    txt(sl, title, Inches(1.1),  y+Inches(0.1),  Inches(11), Inches(0.5), size=17, bold=True, color=col)
    txt(sl, f"문제: {prob}", Inches(1.1), y+Inches(0.6), Inches(5.5), Inches(0.75), size=14, color=C_GRAY)
    txt(sl, f"해결: {sol}",  Inches(6.8), y+Inches(0.6), Inches(5.7), Inches(0.75), size=14)

# ══════════════════════════════════════════════════════════════════════════════
# 15 — 검증 결과
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "Phase 1 검증 결과", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "VALIDATION")

for i, (val, label) in enumerate([
    ("6,055",  "총 청크 수\n(7권 전권)"),
    ("1,024",  "임베딩 벡터\n차원 수"),
    ("~35s",   "1권 인제스트\n소요 시간"),
    ("<1s",    "쿼리 응답\n시간 (검색)"),
]):
    x = Inches(0.55 + i * 3.1)
    card(sl, x, Inches(1.7), Inches(2.8), Inches(1.8))
    txt(sl, val,   x, Inches(1.85), Inches(2.8), Inches(0.8),
        size=38, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    txt(sl, label, x, Inches(2.7),  Inches(2.8), Inches(0.7),
        size=14, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(sl, "컴포넌트 검증 체크리스트", Inches(0.55), Inches(3.7), Inches(12), Inches(0.5),
    size=20, bold=True)
for i, (icon, title, detail) in enumerate([
    ("✅", "PDF 추출",    "238 페이지 정상 추출 (Basic Concepts)"),
    ("✅", "청킹 분류",   "criteria / narrative 정확히 분리"),
    ("✅", "한/영 임베딩","BGE-M3 한국어·영어 동시 벡터화"),
    ("✅", "Qdrant 저장", "6,055 포인트 status=green 확인"),
    ("✅", "XML 파싱",    "AF 진단 + 이상 파라미터 정상 추출"),
    ("✅", "필터 검색",   "language / audience_level 필터 정상 작동"),
]):
    y = Inches(4.3 + i * 0.45)
    txt(sl, icon,  Inches(0.55), y, Inches(0.45), Inches(0.42), size=14)
    txt(sl, title, Inches(1.1),  y, Inches(2.2),  Inches(0.42), size=14, bold=True, color=C_ACCENT)
    txt(sl, detail,Inches(3.4),  y, Inches(9.1),  Inches(0.42), size=14, color=C_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# 16 — Phase 2 로드맵
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "Phase 2 — 품질 향상 로드맵", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "ROADMAP")

for i, (title, desc, benefit) in enumerate([
    ("하이브리드 검색",
     "Dense(의미) + Sparse(키워드)\n동시 검색으로 정밀도 향상",
     "AF ↔ 심방세동 동의어 모두 포착"),
    ("표/그림 캡션 처리",
     "진단 기준 표를 구조화된\n텍스트로 변환 후 별도 태깅",
     '"PR > 200ms" 같은 수치 정확히 추출'),
    ("diagnosis_tags 자동 추출",
     "NER/규칙 기반으로 각 청크에\n진단명 태그 자동 부여",
     "XML 진단명으로 직접 필터 매핑"),
    ("색인 페이지 필터",
     "책 뒷부분 색인 페이지 패턴 감지\n→ 인제스트 시 자동 제외",
     "검색 품질 즉각 향상"),
    ("Cross-Encoder 재순위",
     "벡터 검색 Top-20 → 재순위\n모델로 Top-5 정제",
     "관련도 낮은 결과 제거"),
]):
    row = i // 2
    col = i % 2
    x = Inches(3.7) if i == 4 else Inches(0.55 + col * 6.35)
    y = Inches(1.65 + row * 2.3)
    card(sl, x, y, Inches(5.9), Inches(2.0), border=C_ACCENT2)
    txt(sl, f"{i+1}. {title}", x+Inches(0.2), y+Inches(0.1),  Inches(5.5), Inches(0.5),
        size=16, bold=True, color=C_ACCENT)
    txt(sl, desc,              x+Inches(0.2), y+Inches(0.6),  Inches(5.5), Inches(0.7), size=13)
    txt(sl, f"기대 효과: {benefit}", x+Inches(0.2), y+Inches(1.35), Inches(5.5), Inches(0.5),
        size=12, color=C_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# 17 — Phase 3 확장
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "Phase 3 — 멀티모달 & 확장", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "FUTURE")
txt(sl, "현재: XML(수치) → 미래: 실제 파형 이미지 직접 분석",
    Inches(0.55), Inches(1.55), Inches(12.2), Inches(0.5), size=20, color=C_ACCENT2, bold=True)

for i, (title, desc, tech) in enumerate([
    ("파형 이미지 임베딩",  "ECG 그래프 이미지를 벡터로\n변환해 유사 패턴 검색",  "CLIP / MedCLIP 계열 모델"),
    ("교재 그림 인덱싱",   "교재 속 심전도 그림을 텍스트\n설명과 함께 멀티모달 청크 저장", "그림-텍스트 쌍 검색"),
    ("실시간 스트리밍",    "24시간 홀터 모니터링 데이터\n지속 수집 & 이상 감지 시 RAG", "WebSocket + 이벤트 트리거"),
    ("다국어 확장",        "일본어·중국어 교재 추가\nBGE-M3 그대로 활용 가능", "임베딩 모델 교체 불필요"),
]):
    row = i // 2; col = i % 2
    x = Inches(0.55 + col * 6.35)
    y = Inches(2.3  + row * 2.2)
    card(sl, x, y, Inches(5.9), Inches(1.9), border=C_ORANGE)
    txt(sl, title, x+Inches(0.2), y+Inches(0.1),  Inches(5.5), Inches(0.5), size=16, bold=True, color=C_ORANGE)
    txt(sl, desc,  x+Inches(0.2), y+Inches(0.6),  Inches(5.5), Inches(0.65), size=13)
    txt(sl, f"기술: {tech}", x+Inches(0.2), y+Inches(1.3), Inches(5.5), Inches(0.45), size=12, color=C_ACCENT2)

# ══════════════════════════════════════════════════════════════════════════════
# 18 — 아키텍처 선택 근거
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "아키텍처 선택 근거", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "DESIGN DECISION")

for i, (q, why, benefit) in enumerate([
    ("왜 Qdrant?",
     "로컬 파일 모드로 서버 없이\n바로 시작 가능",
     "서버 모드 전환 시 코드 변경 없음\n하이브리드 검색 내장 지원"),
    ("왜 BGE-M3?",
     "한/영 동시 지원이 필수\n(번역본 포함 7권 커버)",
     "Dense+Sparse 동시 출력\n오픈소스, 외부 API 비용 없음"),
    ("왜 의미 경계 청킹?",
     "진단 기준 수치를 중간에 자르면\n잘못된 기준값이 검색됨",
     "criteria/narrative 분류로\n검색 정밀도 향상"),
    ("왜 메타데이터 필터?",
     "환자 설명 vs 의사 설명이\n달라야 함",
     "audience_level로 검색 소스 분리\n동일 인프라 재사용"),
]):
    row = i // 2; col = i % 2
    x = Inches(0.55 + col * 6.35)
    y = Inches(1.6  + row * 2.5)
    card(sl, x, y, Inches(5.9), Inches(2.3))
    txt(sl, q,       x+Inches(0.2), y+Inches(0.1),  Inches(5.5), Inches(0.5), size=17, bold=True, color=C_ACCENT)
    txt(sl, f"이유: {why}",    x+Inches(0.2), y+Inches(0.6),  Inches(5.5), Inches(0.75), size=13, color=C_GRAY)
    txt(sl, f"장점: {benefit}", x+Inches(0.2), y+Inches(1.4), Inches(5.5), Inches(0.75), size=13, color=C_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# 19 — 프로젝트 구조
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)
accent_bar(sl)
txt(sl, "프로젝트 구조 & 파일 역할", Inches(0.55), Inches(0.7), Inches(10), Inches(0.7),
    size=34, bold=True)
section_tag(sl, "PROJECT STRUCTURE")

card(sl, Inches(0.55), Inches(1.6), Inches(5.5), Inches(5.3))
txt(sl, "디렉토리 구조", Inches(0.8), Inches(1.75), Inches(5.0), Inches(0.45),
    size=16, bold=True, color=C_ACCENT)
bullets(sl, [
    "ecg-rag/",
    "├── data/           # PDF 7권",
    "├── qdrant_storage/ # 벡터 DB 파일",
    "├── config.py       # 전역 설정",
    "├── pipeline.py     # CLI 진입점",
    "└── src/",
    "    ├── ingest/",
    "    │   ├── pdf_extractor.py",
    "    │   ├── chunker.py",
    "    │   └── embedder.py",
    "    ├── db/",
    "    │   └── qdrant_store.py",
    "    └── query/",
    "        ├── xml_parser.py",
    "        └── retriever.py",
], Inches(0.8), Inches(2.3), Inches(5.0), Inches(4.3), size=13, color=C_ACCENT2)

card(sl, Inches(6.4), Inches(1.6), Inches(6.5), Inches(5.3))
txt(sl, "파일별 역할", Inches(6.65), Inches(1.75), Inches(6.0), Inches(0.45),
    size=16, bold=True, color=C_ACCENT)
for i, (fname, role) in enumerate([
    ("pdf_extractor.py", "PDF → 페이지 텍스트"),
    ("chunker.py",       "문장 분리 + 의미 경계 청킹"),
    ("embedder.py",      "BGE-M3 배치 임베딩"),
    ("qdrant_store.py",  "컬렉션 생성·upsert·검색"),
    ("xml_parser.py",    "ECG XML → 쿼리 컨텍스트"),
    ("retriever.py",     "쿼리 임베딩 → 검색 → 포맷"),
    ("pipeline.py",      "ingest / query CLI"),
    ("config.py",        "경로·청크 크기·책 메타"),
]):
    y = Inches(2.3 + i * 0.55)
    txt(sl, fname, Inches(6.65), y, Inches(2.8), Inches(0.5), size=13, color=C_ACCENT2, bold=True, italic=True)
    txt(sl, role,  Inches(9.55), y, Inches(3.2), Inches(0.5), size=13)

# ══════════════════════════════════════════════════════════════════════════════
# 20 — 마무리 & 핵심 요약
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl)

deco = sl.shapes.add_shape(1, Inches(9.5), Inches(-0.3), Inches(4.5), H+Inches(0.3))
deco.fill.solid(); deco.fill.fore_color.rgb = RGBColor(0x08,0x18,0x28); deco.line.fill.background()
line2 = sl.shapes.add_shape(1, Inches(9.47), Inches(0), Pt(3), H)
line2.fill.solid(); line2.fill.fore_color.rgb = C_ACCENT; line2.line.fill.background()

txt(sl, "핵심 요약", Inches(0.55), Inches(0.6), Inches(8.5), Inches(0.65), size=38, bold=True)

for i, (tag, desc) in enumerate([
    ("데이터", "심전도 교재 7권 → 6,055 청크로 분해"),
    ("임베딩", "BGE-M3로 한/영 동일 벡터 공간 구성"),
    ("저장",   "Qdrant + 메타데이터 필터 (레벨·언어)"),
    ("입력",   "ECG XML → 자동 파싱 → 자연어 쿼리"),
    ("검색",   "코사인 유사도 Top-K + 메타데이터 필터"),
    ("출력",   "검색된 교재 근거 기반 LLM 설명 생성"),
]):
    y = Inches(1.5 + i * 0.85)
    tb = sl.shapes.add_shape(1, Inches(0.55), y, Inches(1.4), Inches(0.55))
    tb.fill.solid(); tb.fill.fore_color.rgb = C_ACCENT; tb.line.fill.background()
    tt = tb.text_frame.paragraphs[0]; tt.alignment = PP_ALIGN.CENTER
    tr = tt.add_run(); tr.text = tag; tr.font.size = Pt(14); tr.font.bold = True; tr.font.color.rgb = C_BG
    txt(sl, desc, Inches(2.1), y+Inches(0.06), Inches(7.0), Inches(0.5), size=17)

txt(sl, "Next Steps", Inches(9.8), Inches(1.5), Inches(3.2), Inches(0.5),
    size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
for i, n in enumerate([
    "색인 페이지 필터 추가",
    "하이브리드 검색 적용",
    "diagnosis_tags 자동 태깅",
    "Cross-Encoder 재순위",
    "LLM 연동 (Claude API)",
]):
    y = Inches(2.1 + i * 0.7)
    dot = sl.shapes.add_shape(9, Inches(9.9), y+Inches(0.15), Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = C_ACCENT; dot.line.fill.background()
    txt(sl, n, Inches(10.15), y, Inches(3.0), Inches(0.5), size=14, color=C_GRAY)

out = "/mnt/workspace/ecg-rag/ECG_RAG_발표자료.pptx"
prs.save(out)
print(f"저장 완료: {out}")
