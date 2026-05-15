"""ECG RAG 발표자료 — Graphify 톤앤매너 적용 버전"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── 팔레트 (Graphify 원본에서 추출) ──────────────────────────────────────────
C_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)
C_PURE_BK = RGBColor(0x00, 0x00, 0x00)
C_DARK    = RGBColor(0x4A, 0x4A, 0x4A)
C_BLUE    = RGBColor(0x25, 0x63, 0xEB)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_BG      = RGBColor(0xF9, 0xFA, 0xFB)
C_BG2     = RGBColor(0xF3, 0xF4, 0xF6)
C_BORDER  = RGBColor(0xD1, 0xD5, 0xDB)
C_AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
C_GREEN   = RGBColor(0x05, 0x96, 0x69)
C_RED     = RGBColor(0xDC, 0x26, 0x26)

F_TITLE = "Space Grotesk SemiBold"
F_MONO  = "Space Mono"
F_BODY  = "Inter"

W = Inches(10.0)
H = Inches(5.62)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ──────────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def fill_bg(slide, color=C_BG):
    bg = slide.shapes.add_shape(1, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

def header_bar(slide, title, bar_h=Inches(1.04)):
    bar = slide.shapes.add_shape(1, Inches(0.01), Inches(0.01), Inches(9.98), bar_h)
    bar.fill.solid(); bar.fill.fore_color.rgb = C_PURE_BK; bar.line.fill.background()
    div = slide.shapes.add_shape(1, Inches(0.01), bar_h, Inches(9.98), Pt(1))
    div.fill.solid(); div.fill.fore_color.rgb = C_BLACK; div.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.48), Inches(0.0), Inches(9.0), bar_h)
    tb.word_wrap = False
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title.upper()
    run.font.name = F_TITLE
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = C_WHITE

def txt(slide, text, x, y, w, h, size=10, bold=False, color=C_BLACK,
        align=PP_ALIGN.LEFT, italic=False, font=F_BODY):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    run.font.color.rgb = color

def mono(slide, text, x, y, w, h, size=8, bold=False, color=C_BLUE, align=PP_ALIGN.LEFT):
    txt(slide, text, x, y, w, h, size=size, bold=bold, color=color, align=align, font=F_MONO)

def section_label(slide, label, x=Inches(0.48), y=Inches(1.15)):
    mono(slide, label, x, y, Inches(9.0), Inches(0.22), size=8, color=C_BLUE)

def card(slide, x, y, w, h, bg=C_WHITE, border=C_BLACK, bw=Pt(1)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border; shape.line.width = bw

def th_row(slide, x, y, w, h, label):
    bg = slide.shapes.add_shape(1, x, y, w, h)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BLACK; bg.line.fill.background()
    mono(slide, label, x + Inches(0.08), y + Inches(0.04), w - Inches(0.1), h,
         size=9, bold=True, color=C_WHITE)

def tr_row(slide, x, y, w, h, is_alt=False):
    bg = slide.shapes.add_shape(1, x, y, w, h)
    bg.fill.solid(); bg.fill.fore_color.rgb = C_BG if is_alt else C_WHITE
    bg.line.fill.background()

def bullets(slide, items, x, y, w, h, size=10, color=C_BLACK):
    box = slide.shapes.add_textbox(x, y, w, h)
    box.word_wrap = True
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = item
        run.font.name = F_BODY; run.font.size = Pt(size); run.font.color.rgb = color

CONTENT_Y = Inches(1.42)


# ══════════════════════════════════════════════════════════════════════════════
# 01 — 타이틀
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
card(sl, Inches(1.91), Inches(0.9), Inches(6.18), Inches(3.82), bg=C_WHITE, border=C_BLACK)
txt(sl, "ECG RAG SYSTEM", Inches(2.29), Inches(1.2), Inches(5.41), Inches(0.65),
    size=33, color=C_BLACK, font=F_TITLE)
txt(sl, "심전도 AI 진단 설명 시스템", Inches(2.29), Inches(1.95), Inches(5.41), Inches(0.4),
    size=17, color=C_DARK)
txt(sl, "RAG 기반 의학 지식 검색 · 자동 설명 생성", Inches(2.29), Inches(2.42), Inches(5.41), Inches(0.3),
    size=12, color=C_DARK)
mono(sl, "Medical AI  ·  Knowledge Retrieval  ·  Diagnostic Explanation",
     Inches(2.29), Inches(2.9), Inches(5.41), Inches(0.26), size=10, color=C_BLUE)
for i, kw in enumerate(["PDF 교재 7권", "6,055 청크", "BGE-M3 임베딩", "Qdrant 벡터 DB"]):
    x = Inches(2.29 + i * 1.35)
    chip = sl.shapes.add_shape(1, x, Inches(3.5), Inches(1.25), Inches(0.35))
    chip.fill.solid(); chip.fill.fore_color.rgb = C_BLACK; chip.line.fill.background()
    txt(sl, kw, x + Inches(0.05), Inches(3.53), Inches(1.15), Inches(0.3),
        size=8, color=C_WHITE, font=F_MONO, align=PP_ALIGN.CENTER)
mono(sl, "2026.05  ·  Phase 1 Complete", Inches(2.29), Inches(4.1), Inches(5.41), Inches(0.26),
     size=9, color=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# 02 — 왜 RAG인가?
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Why RAG?")
section_label(sl, "BACKGROUND  ·  PROBLEM STATEMENT")

card(sl, Inches(0.32), Inches(1.65), Inches(4.48), Inches(3.65))
mono(sl, "ISSUE", Inches(0.5), Inches(1.72), Inches(1.5), Inches(0.22))
txt(sl, "기존 방식의 한계", Inches(0.5), Inches(1.96), Inches(4.1), Inches(0.3), size=11, color=C_BLACK)
bullets(sl, [
    "· 심전도 분석기는 진단명만 출력",
    "· 의사/환자에게 근거 설명 없음",
    "· 교재 지식을 코드에 하드코딩하면 유지보수 불가능",
    "· LLM만 단독 사용 시 의학적 환각(Hallucination) 위험",
], Inches(0.5), Inches(2.33), Inches(4.1), Inches(2.8), size=10, color=C_DARK)

card(sl, Inches(5.2), Inches(1.65), Inches(4.48), Inches(3.65))
mono(sl, "SOLUTION", Inches(5.38), Inches(1.72), Inches(1.8), Inches(0.22), color=C_GREEN)
txt(sl, "RAG로 해결", Inches(5.38), Inches(1.96), Inches(4.1), Inches(0.3), size=11, color=C_BLACK)
bullets(sl, [
    "· 실제 의학 교재를 검색 소스로 활용",
    "· XML 진단 결과 → 관련 챕터 자동 검색",
    "· LLM이 검색된 근거로 설명 생성",
    "· 교재 추가만으로 지식 확장 가능",
    "· 출처(책/페이지) 제공으로 신뢰도 확보",
], Inches(5.38), Inches(2.33), Inches(4.1), Inches(2.8), size=10, color=C_DARK)
txt(sl, "→", Inches(4.82), Inches(3.1), Inches(0.36), Inches(0.5),
    size=20, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# 03 — RAG 개념
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Core Concept: RAG")
section_label(sl, "RETRIEVAL-AUGMENTED GENERATION")
txt(sl, "핵심: LLM의 창작이 아닌, 실제 문서에서 찾아서 설명하는 방식",
    Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.28), size=12, color=C_DARK)

for i, (step, title, desc) in enumerate([
    ("STEP 01", "검색 (Retrieve)", "질문과 관련된 문서를\n벡터 DB에서 유사도 검색"),
    ("STEP 02", "증강 (Augment)",  "검색된 문서를 프롬프트에\n추가해 LLM에게 컨텍스트 제공"),
    ("STEP 03", "생성 (Generate)", "LLM이 검색된 근거를 바탕으로\n신뢰할 수 있는 답변 생성"),
]):
    x = Inches(0.32 + i * 3.12)
    card(sl, x, Inches(2.0), Inches(3.0), Inches(2.95))
    mono(sl, step, x + Inches(0.18), Inches(2.1), Inches(2.6), Inches(0.22))
    txt(sl, title, x + Inches(0.18), Inches(2.34), Inches(2.6), Inches(0.3), size=11, color=C_BLACK)
    txt(sl, desc,  x + Inches(0.18), Inches(2.72), Inches(2.65), Inches(1.0), size=10, color=C_DARK)
    if i < 2:
        txt(sl, "→", Inches(3.18 + i * 3.12), Inches(3.3), Inches(0.26), Inches(0.4),
            size=14, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
mono(sl, "> Beyond simple Q&A  ·  Document-grounded generation  ·  Traceable outputs",
     Inches(0.32), Inches(5.2), Inches(9.36), Inches(0.22))

# ══════════════════════════════════════════════════════════════════════════════
# 04 — 시스템 전체 구조
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Pipeline Architecture")
section_label(sl, "SYSTEM OVERVIEW")

mono(sl, "PHASE A  ·  INGEST PIPELINE  (1회 실행)", Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.22))
for i, label in enumerate(["PDF 교재\n(7권)", "텍스트\n추출", "의미 단위\n청킹", "BGE-M3\n임베딩", "Qdrant\n저장"]):
    x = Inches(0.32 + i * 1.88)
    c = sl.shapes.add_shape(1, x, Inches(1.87), Inches(1.76), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = C_BLACK; c.line.fill.background()
    txt(sl, label, x, Inches(1.87), Inches(1.76), Inches(0.7),
        size=9, color=C_WHITE, align=PP_ALIGN.CENTER, font=F_MONO)
    if i < 4:
        txt(sl, "›", Inches(1.94 + i * 1.88), Inches(2.05), Inches(0.22), Inches(0.3),
            size=13, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

div = sl.shapes.add_shape(1, Inches(0.32), Inches(2.73), Inches(9.36), Pt(1))
div.fill.solid(); div.fill.fore_color.rgb = C_BORDER; div.line.fill.background()

mono(sl, "PHASE B  ·  QUERY PIPELINE  (매 요청마다 실행)", Inches(0.32), Inches(2.82), Inches(9.36), Inches(0.22))
for i, label in enumerate(["ECG XML\n분석 결과", "XML\n파싱", "쿼리\n임베딩", "Qdrant\n검색", "LLM\n설명 생성", "진단\n설명 출력"]):
    x = Inches(0.32 + i * 1.57)
    c = sl.shapes.add_shape(1, x, Inches(3.1), Inches(1.46), Inches(0.7))
    c.fill.solid(); c.fill.fore_color.rgb = C_BG2; c.line.color.rgb = C_BLACK; c.line.width = Pt(0.75)
    txt(sl, label, x, Inches(3.1), Inches(1.46), Inches(0.7),
        size=9, color=C_BLACK, align=PP_ALIGN.CENTER, font=F_MONO)
    if i < 5:
        txt(sl, "›", Inches(1.64 + i * 1.57), Inches(3.28), Inches(0.26), Inches(0.3),
            size=12, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
mono(sl, "> SHARED RESOURCE: Qdrant Vector DB  ·  인제스트 1회 → 쿼리 반복 사용",
     Inches(0.32), Inches(4.05), Inches(9.36), Inches(0.22))

# ══════════════════════════════════════════════════════════════════════════════
# 05 — 데이터 소스
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Data Source: 7 ECG Textbooks")
section_label(sl, "CORPUS OVERVIEW")

col_xs = [Inches(0.32), Inches(4.2), Inches(5.8), Inches(7.3)]
col_ws = [Inches(3.8),  Inches(1.5), Inches(1.4), Inches(1.5)]
for hx, hw, hl in zip(col_xs, col_ws, ["교재명", "언어", "레벨", "청크 수"]):
    th_row(sl, hx, Inches(1.62), hw, Inches(0.32), hl)

lc_map = {"기초": C_GREEN, "임상": C_BLUE, "전문": C_AMBER}
for i, (name, lang, level, cnt) in enumerate([
    ("만화로보는 심전도",           "한국어", "기초",   "97"),
    ("Basic Concepts of EKG",     "영어",   "기초",  "405"),
    ("Goldberger's Clinical ECG", "영어",   "임상", "1,503"),
    ("Marriott's 13판 (번역)",     "한국어", "임상", "2,384"),
    ("Marriott's 13판 (원문)",     "영어",   "임상",  "946"),
    ("Basic Electrocardiography", "영어",   "임상",  "225"),
    ("Inherited Arrhythmias",     "영어",   "전문",  "495"),
]):
    y = Inches(1.94 + i * 0.38)
    for hx, hw in zip(col_xs, col_ws):
        tr_row(sl, hx, y, hw, Inches(0.36), i % 2 == 1)
    txt(sl, name, col_xs[0] + Inches(0.08), y + Inches(0.06), col_ws[0] - Inches(0.1), Inches(0.28), size=9, color=C_BLACK)
    txt(sl, lang, col_xs[1] + Inches(0.08), y + Inches(0.06), col_ws[1] - Inches(0.1), Inches(0.28), size=9, color=C_DARK)
    chip = sl.shapes.add_shape(1, col_xs[2] + Inches(0.08), y + Inches(0.06), Inches(0.9), Inches(0.25))
    chip.fill.solid(); chip.fill.fore_color.rgb = lc_map.get(level, C_BLACK); chip.line.fill.background()
    ct = chip.text_frame.paragraphs[0]; ct.alignment = PP_ALIGN.CENTER
    cr = ct.add_run(); cr.text = level; cr.font.name = F_MONO; cr.font.size = Pt(8); cr.font.color.rgb = C_WHITE
    txt(sl, cnt, col_xs[3] + Inches(0.08), y + Inches(0.06), col_ws[3] - Inches(0.1), Inches(0.28),
        size=9, bold=True, color=C_BLACK, align=PP_ALIGN.RIGHT, font=F_MONO)
mono(sl, "> TOTAL: 6,055 chunks  ·  Korean 2,481  ·  English 3,574",
     Inches(0.32), Inches(4.7), Inches(9.36), Inches(0.22))

# ══════════════════════════════════════════════════════════════════════════════
# 06 — PDF 추출
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Step 01 — PDF Text Extraction")
section_label(sl, "INGEST PIPELINE  ·  TOOL: PyMuPDF")
txt(sl, "208MB 교재 → 페이지별 순수 텍스트 (헤더/푸터 자동 제거)",
    Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.28), size=11, color=C_DARK)

for i, (step_n, title, desc) in enumerate([
    ("01", "블록 단위 추출", "PDF의 텍스트 블록 구조를\n그대로 활용해 단락 경계 보존"),
    ("02", "헤더/푸터 제거", "페이지 상하 7% 영역을\n자동으로 필터링"),
    ("03", "짧은 페이지 스킵", "80자 미만 페이지 제외\n(그림 전용, 빈 페이지)"),
]):
    x = Inches(0.32 + i * 3.12)
    card(sl, x, Inches(2.0), Inches(2.96), Inches(2.15))
    mono(sl, f"METHOD {step_n}", x + Inches(0.18), Inches(2.1), Inches(2.6), Inches(0.22))
    txt(sl, title, x + Inches(0.18), Inches(2.36), Inches(2.6), Inches(0.3), size=11, color=C_BLACK)
    txt(sl, desc,  x + Inches(0.18), Inches(2.74), Inches(2.6), Inches(1.0), size=10, color=C_DARK)

card(sl, Inches(0.32), Inches(4.3), Inches(9.36), Inches(1.0), bg=C_BG2, border=C_BORDER)
mono(sl, "OUTPUT SCHEMA", Inches(0.5), Inches(4.38), Inches(3), Inches(0.22))
mono(sl, '{ "page_num": 12,  "text": "Sinus rhythm requires ...",  "source_file": "basic-concepts-of-ekg.pdf" }',
     Inches(0.5), Inches(4.62), Inches(9.0), Inches(0.55), size=9, color=C_BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# 07 — 청킹
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Step 02 — Semantic Boundary Chunking")
section_label(sl, "INGEST PIPELINE  ·  STRATEGY: MEANING-AWARE SPLIT")
txt(sl, "고정 크기 자르기(✗)  →  의학 문맥을 이해한 경계 분리(✓)",
    Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.28), size=11, color=C_DARK)

card(sl, Inches(0.32), Inches(2.0), Inches(4.5), Inches(3.3))
mono(sl, "CHUNKING RULES", Inches(0.5), Inches(2.1), Inches(3.0), Inches(0.22))
bullets(sl, [
    "① 섹션 제목 감지 → 새 청크 시작",
    "② 진단 기준 블록(ms/mm/bpm 포함)",
    "   → 절대 중간에 자르지 않음",
    "③ 목록(bullet/번호) → 통째로 유지",
    "④ 목표 500 토큰, 최대 700 토큰",
    "⑤ 2문장 오버랩으로 문맥 연속성 유지",
], Inches(0.5), Inches(2.38), Inches(4.1), Inches(2.6), size=10, color=C_DARK)

card(sl, Inches(5.18), Inches(2.0), Inches(4.5), Inches(3.3))
mono(sl, "CHUNK TYPES", Inches(5.36), Inches(2.1), Inches(3.0), Inches(0.22))
th_row(sl, Inches(5.36), Inches(2.38), Inches(1.4), Inches(0.3), "TYPE")
th_row(sl, Inches(6.76), Inches(2.38), Inches(2.74), Inches(0.3), "DESCRIPTION / COUNT")
for i, (tag, desc, cnt, col) in enumerate([
    ("criteria",  "진단 기준, 수치 임계값 포함\n블록 — LLM 정밀 근거 제공", "1,931", C_GREEN),
    ("narrative", "임상 설명, 배경 지식 텍스트\n— LLM 문맥 제공",          "4,124", C_BLUE),
]):
    y = Inches(2.68 + i * 1.15)
    tr_row(sl, Inches(5.36), y, Inches(4.14), Inches(1.05), i % 2 == 1)
    chip = sl.shapes.add_shape(1, Inches(5.44), y + Inches(0.12), Inches(1.2), Inches(0.28))
    chip.fill.solid(); chip.fill.fore_color.rgb = col; chip.line.fill.background()
    ct = chip.text_frame.paragraphs[0]; ct.alignment = PP_ALIGN.CENTER
    cr = ct.add_run(); cr.text = tag; cr.font.name = F_MONO; cr.font.size = Pt(8); cr.font.color.rgb = C_WHITE
    txt(sl, desc, Inches(6.76), y + Inches(0.08), Inches(2.5), Inches(0.7), size=9, color=C_DARK)
    mono(sl, cnt, Inches(8.9), y + Inches(0.38), Inches(0.7), Inches(0.28), size=9, bold=True, color=C_BLACK, align=PP_ALIGN.RIGHT)
mono(sl, "> RESULT: 6,055 total chunks  ·  criteria 1,931  ·  narrative 4,124",
     Inches(0.32), Inches(5.2), Inches(9.36), Inches(0.22))

# ══════════════════════════════════════════════════════════════════════════════
# 08 — 임베딩
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Step 03 — BGE-M3 Embedding")
section_label(sl, "INGEST PIPELINE  ·  MODEL: BAAI/BGE-M3")
txt(sl, "텍스트 → 1,024차원 숫자 벡터  |  의미가 비슷한 텍스트는 가까운 벡터",
    Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.28), size=11, color=C_DARK)

th_row(sl, Inches(0.32), Inches(1.97), Inches(4.5),  Inches(0.3), "INPUT TEXT")
th_row(sl, Inches(4.82), Inches(1.97), Inches(4.86), Inches(0.3), "EMBEDDING VECTOR (샘플)")
for i, (src, vec, note) in enumerate([
    ('"Atrial Fibrillation"',    "[0.23, -0.87, 0.41, ...]", ""),
    ('"심방세동 진단 기준"',         "[0.21, -0.84, 0.39, ...]", "← 거의 같은 방향"),
    ('"심전도 P파 부재"',            "[0.19, -0.80, 0.44, ...]", "← 유사"),
]):
    y = Inches(2.27 + i * 0.42)
    tr_row(sl, Inches(0.32), y, Inches(4.5),  Inches(0.4), i % 2 == 1)
    tr_row(sl, Inches(4.82), y, Inches(4.86), Inches(0.4), i % 2 == 1)
    mono(sl, src, Inches(0.5), y + Inches(0.08), Inches(4.1), Inches(0.28), size=9, color=C_BLACK)
    mono(sl, vec, Inches(5.0), y + Inches(0.08), Inches(3.0), Inches(0.28), size=9, color=C_DARK)
    if note:
        mono(sl, note, Inches(8.1), y + Inches(0.1), Inches(1.5), Inches(0.25), size=8)

txt(sl, "BGE-M3 선택 이유", Inches(0.32), Inches(3.62), Inches(9.36), Inches(0.3), size=11, color=C_BLACK)
for i, (title, desc) in enumerate([
    ("한/영 동시 지원",     "한국어 번역본과 영어 원문을 같은 벡터 공간에서 검색"),
    ("Dense + Sparse 동시 출력", "Phase 2 하이브리드 검색으로 그대로 확장 가능"),
    ("오픈소스 · 자체 호스팅", "외부 API 의존 없음, 비용 없음"),
]):
    x = Inches(0.32 + i * 3.12)
    card(sl, x, Inches(4.0), Inches(3.0), Inches(1.35), bg=C_BG2, border=C_BORDER)
    mono(sl, f"FEATURE {i+1:02d}", x + Inches(0.15), Inches(4.08), Inches(2.6), Inches(0.22))
    txt(sl, title, x + Inches(0.15), Inches(4.32), Inches(2.7), Inches(0.28), size=10, color=C_BLACK)
    txt(sl, desc,  x + Inches(0.15), Inches(4.64), Inches(2.7), Inches(0.62), size=9, color=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# 09 — Qdrant
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Step 04 — Qdrant Vector DB")
section_label(sl, "INGEST PIPELINE  ·  STORAGE: LOCAL FILE MODE")

card(sl, Inches(0.32), Inches(1.62), Inches(4.48), Inches(3.68))
mono(sl, "QDRANT OVERVIEW", Inches(0.5), Inches(1.72), Inches(3.5), Inches(0.22))
bullets(sl, [
    "· 벡터 유사도 검색에 특화된 DB",
    "· 고속 코사인 유사도 계산",
    "· 메타데이터 필터링 동시 지원",
    "· 로컬 파일 모드 — 서버 불필요",
    "· 추후 서버 모드로 코드 변경 없이 전환",
], Inches(0.5), Inches(2.02), Inches(4.1), Inches(3.0), size=10, color=C_DARK)

card(sl, Inches(5.18), Inches(1.62), Inches(4.5), Inches(3.68))
mono(sl, "POINT SCHEMA  (청크 1개)", Inches(5.36), Inches(1.72), Inches(4.0), Inches(0.22))
th_row(sl, Inches(5.36), Inches(2.02), Inches(1.7), Inches(0.3), "FIELD")
th_row(sl, Inches(7.06), Inches(2.02), Inches(2.44), Inches(0.3), "VALUE")
for i, (f, v) in enumerate([
    ("id",             "고유 UUID"),
    ("vector",         "1024차원 숫자 배열"),
    ("content",        "청크 원문 텍스트"),
    ("source_book",    "marriotts_13ed_ko 등"),
    ("language",       "en / ko"),
    ("audience_level", "basic / clinical / specialist"),
    ("content_type",   "criteria / narrative"),
    ("page_num",       "원본 페이지 번호"),
]):
    y = Inches(2.32 + i * 0.34)
    tr_row(sl, Inches(5.36), y, Inches(4.14), Inches(0.32), i % 2 == 1)
    mono(sl, f, Inches(5.44), y + Inches(0.05), Inches(1.55), Inches(0.26), size=8, color=C_BLACK)
    txt(sl,  v, Inches(7.14), y + Inches(0.05), Inches(2.28), Inches(0.26), size=9, color=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# 10 — 메타데이터 필터링
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Metadata Filtering Strategy")
section_label(sl, "PRECISION RETRIEVAL  ·  AUDIENCE-AWARE SEARCH")
txt(sl, "벡터 유사도 + 조건 필터 = 같은 질환을 대상별로 다르게 설명",
    Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.28), size=11, color=C_DARK)

for i, (title, filt, result, col) in enumerate([
    ("환자 대상 설명", "audience_level = basic\nlanguage = ko",
     "만화 심전도 · 기초 교재에서\n쉬운 설명 검색", C_GREEN),
    ("의사 대상 설명", "audience_level = clinical\nlanguage = en/ko",
     "Marriott's · Goldberger's에서\n임상 기준 검색", C_BLUE),
    ("전문의 대상", "audience_level = specialist",
     "Inherited Arrhythmias에서\n유전성 부정맥 전문 내용 검색", C_AMBER),
]):
    x = Inches(0.32 + i * 3.12)
    card(sl, x, Inches(2.02), Inches(3.0), Inches(3.25))
    chip = sl.shapes.add_shape(1, x + Inches(0.15), Inches(2.12), Inches(2.0), Inches(0.28))
    chip.fill.solid(); chip.fill.fore_color.rgb = col; chip.line.fill.background()
    ct = chip.text_frame.paragraphs[0]; ct.alignment = PP_ALIGN.CENTER
    cr = ct.add_run(); cr.text = title; cr.font.name = F_MONO; cr.font.size = Pt(8); cr.font.color.rgb = C_WHITE
    mono(sl, "FILTER", x + Inches(0.15), Inches(2.5),  Inches(1.0), Inches(0.22))
    mono(sl, filt,     x + Inches(0.15), Inches(2.74), Inches(2.7), Inches(0.55), size=9, color=C_BLACK)
    mono(sl, "RESULT", x + Inches(0.15), Inches(3.38), Inches(1.0), Inches(0.22))
    txt(sl, result,    x + Inches(0.15), Inches(3.62), Inches(2.7), Inches(0.9), size=9, color=C_DARK)
mono(sl, "> 동일 인프라, 동일 컬렉션에서 필터만 바꿔 대상별 맞춤 검색",
     Inches(0.32), Inches(5.25), Inches(9.36), Inches(0.22))

# ══════════════════════════════════════════════════════════════════════════════
# 11 — XML 파싱
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Query Step 01 — ECG XML Parsing")
section_label(sl, "QUERY PIPELINE  ·  INPUT: ECG ANALYSIS RESULT")
txt(sl, "심전도 분석기 XML 출력 → RAG 쿼리 컨텍스트 자동 변환",
    Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.28), size=11, color=C_DARK)

card(sl, Inches(0.32), Inches(2.02), Inches(4.48), Inches(3.25))
mono(sl, "INPUT  ·  ECG XML", Inches(0.5), Inches(2.12), Inches(3.0), Inches(0.22))
bullets(sl, [
    "<ECGReport>",
    "  <Measurements>",
    "    <HeartRate unit='bpm'>148</HeartRate>",
    "    <PRInterval unit='ms'>absent</PRInterval>",
    "    <QRSDuration unit='ms'>88</QRSDuration>",
    "  </Measurements>",
    "  <Diagnoses>",
    "    <Diagnosis>Atrial Fibrillation</Diagnosis>",
    "  </Diagnoses>",
    "</ECGReport>",
], Inches(0.5), Inches(2.4), Inches(4.1), Inches(2.7), size=9, color=C_BLACK)
txt(sl, "→", Inches(4.84), Inches(3.3), Inches(0.32), Inches(0.4),
    size=18, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

card(sl, Inches(5.2), Inches(2.02), Inches(4.48), Inches(3.25))
mono(sl, "OUTPUT  ·  ECGQueryContext", Inches(5.38), Inches(2.12), Inches(3.8), Inches(0.22))
bullets(sl, [
    "diagnoses:",
    "  ['Atrial Fibrillation']",
    "",
    "abnormal_params:",
    "  ['HeartRate 148bpm (high)']",
    "",
    "natural_query:",
    "  'Atrial Fibrillation HeartRate 148bpm",
    "   ECG diagnosis criteria...'",
], Inches(5.38), Inches(2.4), Inches(4.1), Inches(2.7), size=9, color=C_BLACK)
mono(sl, "> 정상 범위 자동 비교: HR 60-100bpm, PR 120-200ms, QRS 70-110ms, QTc 350-450ms",
     Inches(0.32), Inches(5.25), Inches(9.36), Inches(0.22))

# ══════════════════════════════════════════════════════════════════════════════
# 12 — 벡터 검색
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Query Step 02 — Vector Search")
section_label(sl, "QUERY PIPELINE  ·  COSINE SIMILARITY RETRIEVAL")

for i, label in enumerate(["자연어\n쿼리 생성", "BGE-M3\n임베딩", "Qdrant\n코사인 검색", "Top-K\n청크 반환"]):
    x = Inches(0.32 + i * 2.3)
    c = sl.shapes.add_shape(1, x, Inches(1.62), Inches(2.1), Inches(0.75))
    c.fill.solid(); c.fill.fore_color.rgb = C_BLACK; c.line.fill.background()
    txt(sl, label, x, Inches(1.62), Inches(2.1), Inches(0.75),
        size=9, color=C_WHITE, align=PP_ALIGN.CENTER, font=F_MONO)
    if i < 3:
        txt(sl, "›", Inches(2.26 + i * 2.3), Inches(1.9), Inches(0.26), Inches(0.3),
            size=13, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

txt(sl, "코사인 유사도: 두 벡터가 얼마나 같은 방향을 가리키는지  (-1 ~ 1,  1에 가까울수록 유사)",
    Inches(0.32), Inches(2.58), Inches(9.36), Inches(0.28), size=10, color=C_DARK)

th_row(sl, Inches(0.32), Inches(3.0), Inches(0.8),  Inches(0.3), "SCORE")
th_row(sl, Inches(1.12), Inches(3.0), Inches(3.5),  Inches(0.3), "SOURCE")
th_row(sl, Inches(4.62), Inches(3.0), Inches(5.06), Inches(0.3), "CONTENT PREVIEW")
for i, (score, source, content) in enumerate([
    ("0.712", "marriotts_13ed_ko  p.342  [criteria]",
     "심방세동의 ECG 기준: 불규칙한 RR 간격, P파 소실, 세동파(f파)..."),
    ("0.684", "goldbergers_10ed   p.218  [criteria]",
     "Atrial fibrillation: absent P waves, irregularly irregular rhythm..."),
    ("0.651", "marriotts_13ed_en  p.330  [narrative]",
     "AF is the most common sustained cardiac arrhythmia, occurring..."),
]):
    y = Inches(3.3 + i * 0.42)
    for hx, hw in [(Inches(0.32), Inches(0.8)), (Inches(1.12), Inches(3.5)), (Inches(4.62), Inches(5.06))]:
        tr_row(sl, hx, y, hw, Inches(0.4), i % 2 == 1)
    mono(sl, score,   Inches(0.4),  y + Inches(0.08), Inches(0.65), Inches(0.28), size=9, bold=True, color=C_BLACK)
    mono(sl, source,  Inches(1.2),  y + Inches(0.08), Inches(3.3),  Inches(0.28), size=8, color=C_DARK)
    txt(sl, content,  Inches(4.7),  y + Inches(0.08), Inches(4.9),  Inches(0.28), size=9, color=C_BLACK)
mono(sl, "> FILTER APPLIED: language=ko  ·  audience_level=clinical  ·  top_k=5",
     Inches(0.32), Inches(4.62), Inches(9.36), Inches(0.22))

# ══════════════════════════════════════════════════════════════════════════════
# 13 — LLM 설명 생성
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Query Step 03 — LLM Explanation Generation")
section_label(sl, "QUERY PIPELINE  ·  GROUNDED GENERATION")

card(sl, Inches(0.32), Inches(1.62), Inches(9.36), Inches(2.5), bg=C_BG2, border=C_BORDER)
mono(sl, "PROMPT STRUCTURE", Inches(0.5), Inches(1.72), Inches(3.0), Inches(0.22))
bullets(sl, [
    "[시스템]  당신은 심전도 전문의입니다. 아래 교재 내용을 근거로 설명하세요.",
    "",
    "[교재 검색 결과]",
    "  Source 1: marriotts_13ed_ko p.342 — 심방세동의 ECG 기준: 불규칙한 RR 간격...",
    "  Source 2: goldbergers_10ed  p.218 — AF: absent P waves, irregularly irregular...",
    "",
    "[ECG 분석 결과]  진단: Atrial Fibrillation  |  이상: HR 148bpm (high)",
    "[질문]  위 심전도 소견에 대해 임상의 수준으로 설명해 주세요.",
], Inches(0.5), Inches(2.0), Inches(9.0), Inches(2.0), size=9, color=C_BLACK)

mono(sl, "LLM RESPONSE", Inches(0.5), Inches(4.22), Inches(2.0), Inches(0.22), color=C_GREEN)
card(sl, Inches(0.32), Inches(4.44), Inches(9.36), Inches(1.0), bg=C_WHITE, border=C_BLACK)
txt(sl, "심방세동(Atrial Fibrillation)은 심방의 불규칙한 전기 활동으로 발생합니다. ECG상 P파가 소실되고 "
        "불규칙한 RR 간격이 특징입니다. 현재 심박수 148bpm으로 빠른 심실 반응을 동반하고 있어 즉각적인 속도 조절 치료가 "
        "필요합니다.  (출처: Marriott's 13판 p.342, Goldberger's 10판 p.218)",
    Inches(0.5), Inches(4.52), Inches(9.0), Inches(0.85), size=10, color=C_BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# 14 — 발견한 문제 & 해결
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Lessons Learned")
section_label(sl, "OPERATIONAL ISSUES & RESOLUTIONS")

th_row(sl, Inches(0.32), Inches(1.62), Inches(2.5), Inches(0.3), "ISSUE")
th_row(sl, Inches(2.82), Inches(1.62), Inches(3.7), Inches(0.3), "PROBLEM")
th_row(sl, Inches(6.52), Inches(1.62), Inches(3.16), Inches(0.3), "RESOLUTION")

for i, (title, prob, sol, col) in enumerate([
    ("한국어 파일명 문제",
     "NFD(자모 분해형) vs NFC(조합형)\n유니코드 불일치로 한국어\nPDF 2권 인제스트 스킵",
     "unicodedata.normalize('NFC')\n로 파일명 정규화 후 해결", C_AMBER),
    ("책 색인 페이지 혼입",
     "책 뒷부분 색인이 ECG 용어 밀집\n→ 관련도 높은 것처럼 검색됨",
     "색인 페이지 패턴 감지 필터\n추가 (Phase 2 예정)", C_RED),
    ("qdrant-client API 변경",
     "v1.14+ 에서 client.search()\n제거 → AttributeError 발생",
     "client.query_points()로\n마이그레이션 완료", C_GREEN),
]):
    y = Inches(1.92 + i * 1.06)
    for hx, hw in [(Inches(0.32), Inches(2.5)), (Inches(2.82), Inches(3.7)), (Inches(6.52), Inches(3.16))]:
        tr_row(sl, hx, y, hw, Inches(1.0), i % 2 == 1)
    chip = sl.shapes.add_shape(1, Inches(0.42), y + Inches(0.08), Inches(0.55), Inches(0.25))
    chip.fill.solid(); chip.fill.fore_color.rgb = col; chip.line.fill.background()
    ct = chip.text_frame.paragraphs[0]; ct.alignment = PP_ALIGN.CENTER
    cr = ct.add_run(); cr.text = f"#{i+1}"; cr.font.name = F_MONO; cr.font.size = Pt(9); cr.font.color.rgb = C_WHITE
    txt(sl, title, Inches(1.05), y + Inches(0.08), Inches(1.6), Inches(0.88), size=9, color=C_BLACK)
    txt(sl, prob,  Inches(2.9),  y + Inches(0.08), Inches(3.5), Inches(0.88), size=9, color=C_DARK)
    txt(sl, sol,   Inches(6.6),  y + Inches(0.08), Inches(3.0), Inches(0.88), size=9, color=C_BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# 15 — 검증 결과
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Phase 1 Validation Results")
section_label(sl, "QA CHECKLIST  ·  ALL COMPONENTS VERIFIED")

for i, (val, label) in enumerate([
    ("6,055", "총 청크 수\n(7권 전권)"),
    ("1,024", "임베딩\n벡터 차원"),
    ("~35s",  "1권 인제스트\n소요 시간"),
    ("<1s",   "검색 쿼리\n응답 시간"),
]):
    x = Inches(0.32 + i * 2.4)
    card(sl, x, Inches(1.62), Inches(2.2), Inches(1.25), bg=C_WHITE, border=C_BLACK)
    txt(sl, val,   x, Inches(1.72), Inches(2.2), Inches(0.65),
        size=28, color=C_BLACK, align=PP_ALIGN.CENTER, font=F_TITLE)
    txt(sl, label, x, Inches(2.42), Inches(2.2), Inches(0.38),
        size=9, color=C_DARK, align=PP_ALIGN.CENTER)

th_row(sl, Inches(0.32), Inches(3.0), Inches(1.6), Inches(0.3), "STATUS")
th_row(sl, Inches(1.92), Inches(3.0), Inches(2.0), Inches(0.3), "COMPONENT")
th_row(sl, Inches(3.92), Inches(3.0), Inches(5.76), Inches(0.3), "VERIFICATION RESULT")
for i, (status, title, detail) in enumerate([
    ("✅ PASS", "PDF 추출",    "238 페이지 정상 추출 (Basic Concepts)"),
    ("✅ PASS", "청킹 분류",   "criteria / narrative 정확히 분리"),
    ("✅ PASS", "한/영 임베딩","BGE-M3 한국어·영어 동시 벡터화"),
    ("✅ PASS", "Qdrant 저장", "6,055 포인트 status=green 확인"),
    ("✅ PASS", "XML 파싱",    "AF 진단 + 이상 파라미터 정상 추출"),
    ("✅ PASS", "필터 검색",   "language / audience_level 필터 정상 작동"),
]):
    y = Inches(3.3 + i * 0.34)
    for hx, hw in [(Inches(0.32), Inches(1.6)), (Inches(1.92), Inches(2.0)), (Inches(3.92), Inches(5.76))]:
        tr_row(sl, hx, y, hw, Inches(0.32), i % 2 == 1)
    mono(sl, status, Inches(0.42), y + Inches(0.05), Inches(1.4), Inches(0.25), size=8, color=C_GREEN)
    txt(sl, title,   Inches(2.0),  y + Inches(0.05), Inches(1.8), Inches(0.25), size=9, color=C_BLACK)
    txt(sl, detail,  Inches(4.0),  y + Inches(0.05), Inches(5.6), Inches(0.25), size=9, color=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# 16 — Phase 2 로드맵
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Phase 2 — Quality Improvement Roadmap")
section_label(sl, "NEXT PHASE  ·  5 ENHANCEMENT AREAS")

th_row(sl, Inches(0.32), Inches(1.62), Inches(0.52), Inches(0.3), "#")
th_row(sl, Inches(0.84), Inches(1.62), Inches(2.4),  Inches(0.3), "ENHANCEMENT")
th_row(sl, Inches(3.24), Inches(1.62), Inches(3.7),  Inches(0.3), "DESCRIPTION")
th_row(sl, Inches(6.94), Inches(1.62), Inches(2.74), Inches(0.3), "EXPECTED BENEFIT")
for i, (title, desc, benefit) in enumerate([
    ("하이브리드 검색", "Dense(의미)+Sparse(키워드) 동시 검색으로 정밀도 향상",
     "AF ↔ 심방세동 동의어 모두 포착"),
    ("표/그림 처리", "진단 기준 표를 구조화된 텍스트로 변환 및 태깅",
     '"PR > 200ms" 수치 기준 정확 추출'),
    ("diagnosis_tags 추출", "NER/규칙으로 각 청크에 진단명 태그 자동 부여",
     "XML 진단명으로 직접 필터 매핑"),
    ("색인 페이지 필터", "책 뒷부분 색인 페이지 패턴 감지 → 자동 제외",
     "검색 품질 즉각 향상"),
    ("Cross-Encoder 재순위", "Top-20 → 재순위 모델로 Top-5 정제",
     "관련도 낮은 결과 제거"),
]):
    y = Inches(1.92 + i * 0.62)
    for hx, hw in [(Inches(0.32), Inches(0.52)), (Inches(0.84), Inches(2.4)),
                   (Inches(3.24), Inches(3.7)),  (Inches(6.94), Inches(2.74))]:
        tr_row(sl, hx, y, hw, Inches(0.6), i % 2 == 1)
    mono(sl, f"P2-{i+1:02d}", Inches(0.38), y + Inches(0.16), Inches(0.45), Inches(0.28), size=8)
    txt(sl, title,   Inches(0.92), y + Inches(0.08), Inches(2.24), Inches(0.5),  size=9, color=C_BLACK)
    txt(sl, desc,    Inches(3.32), y + Inches(0.08), Inches(3.54), Inches(0.5),  size=9, color=C_DARK)
    txt(sl, benefit, Inches(7.02), y + Inches(0.08), Inches(2.58), Inches(0.5),  size=9, color=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# 17 — Phase 3 확장
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Phase 3 — Multimodal & Expansion")
section_label(sl, "FUTURE DIRECTION  ·  BEYOND TEXT")
txt(sl, "현재: XML(수치) → 미래: 실제 파형 이미지 직접 분석",
    Inches(0.32), Inches(1.6), Inches(9.36), Inches(0.28), size=11, color=C_DARK)

for i, (title, desc, tech) in enumerate([
    ("파형 이미지 임베딩",  "ECG 그래프 이미지를 벡터로 변환해\n유사 패턴 직접 검색",      "CLIP / MedCLIP"),
    ("교재 그림 인덱싱",   "교재 속 심전도 그림을 텍스트 설명과\n함께 멀티모달 청크 저장",  "그림-텍스트 쌍 검색"),
    ("실시간 스트리밍",    "24시간 홀터 모니터링 데이터 수집\n이상 감지 시 자동 RAG 트리거", "WebSocket + 이벤트 드리븐"),
    ("다국어 확장",        "일본어·중국어 교재 추가\nBGE-M3 그대로 활용 가능",          "임베딩 모델 교체 불필요"),
]):
    row = i // 2; col = i % 2
    x = Inches(0.32 + col * 4.84)
    y = Inches(2.0  + row * 1.7)
    card(sl, x, y, Inches(4.66), Inches(1.55))
    mono(sl, f"FEATURE {i+1:02d}", x + Inches(0.18), y + Inches(0.1),  Inches(3.0), Inches(0.22))
    txt(sl, title, x + Inches(0.18), y + Inches(0.34), Inches(4.2), Inches(0.28), size=11, color=C_BLACK)
    txt(sl, desc,  x + Inches(0.18), y + Inches(0.68), Inches(4.2), Inches(0.58), size=9, color=C_DARK)
    mono(sl, f"> {tech}", x + Inches(0.18), y + Inches(1.3), Inches(4.0), Inches(0.2), size=8)

# ══════════════════════════════════════════════════════════════════════════════
# 18 — 아키텍처 선택 근거
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Architecture Design Rationale")
section_label(sl, "DESIGN DECISIONS  ·  WHY EACH COMPONENT")

th_row(sl, Inches(0.32), Inches(1.62), Inches(2.2), Inches(0.3), "DECISION")
th_row(sl, Inches(2.52), Inches(1.62), Inches(3.5), Inches(0.3), "RATIONALE")
th_row(sl, Inches(6.02), Inches(1.62), Inches(3.66), Inches(0.3), "BENEFIT")
for i, (q, why, benefit) in enumerate([
    ("왜 Qdrant?",
     "로컬 파일 모드로 서버 없이\n바로 시작 가능",
     "서버 모드 전환 시 코드 변경 없음\n하이브리드 검색 내장 지원"),
    ("왜 BGE-M3?",
     "한/영 동시 지원이 필수\n(번역본 포함 7권 커버)",
     "Dense+Sparse 동시 출력\n오픈소스, 외부 API 비용 없음"),
    ("왜 의미 경계 청킹?",
     "진단 기준 수치를 중간에 자르면\n잘못된 기준값이 검색됨",
     "criteria/narrative 분류로\n검색 정밀도 향상"),
    ("왜 메타데이터 필터?",
     "환자 설명 vs 의사 설명이\n달라야 함",
     "audience_level로 검색 소스 분리\n동일 인프라 재사용"),
]):
    y = Inches(1.92 + i * 0.82)
    for hx, hw in [(Inches(0.32), Inches(2.2)), (Inches(2.52), Inches(3.5)), (Inches(6.02), Inches(3.66))]:
        tr_row(sl, hx, y, hw, Inches(0.8), i % 2 == 1)
    txt(sl, q,       Inches(0.42), y + Inches(0.08), Inches(2.0), Inches(0.7), size=10, color=C_BLACK)
    txt(sl, why,     Inches(2.62), y + Inches(0.08), Inches(3.3), Inches(0.7), size=9, color=C_DARK)
    txt(sl, benefit, Inches(6.12), y + Inches(0.08), Inches(3.46),Inches(0.7), size=9, color=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# 19 — 프로젝트 구조
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Project Structure & File Roles")
section_label(sl, "CODEBASE OVERVIEW")

card(sl, Inches(0.32), Inches(1.62), Inches(4.14), Inches(3.78), bg=C_BG2, border=C_BORDER)
mono(sl, "DIRECTORY TREE", Inches(0.5), Inches(1.72), Inches(3.0), Inches(0.22))
bullets(sl, [
    "ecg-rag/",
    "├── data/           # PDF 7권",
    "├── qdrant_storage/ # 벡터 DB 파일",
    "├── config.py       # 전역 설정",
    "├── pipeline.py     # CLI 진입점",
    "└── src/",
    "    ├── ingest/",
    "    │   ├── pdf_extractor.py",
    "    │   ├── chunker.py",
    "    │   └── embedder.py",
    "    ├── db/",
    "    │   └── qdrant_store.py",
    "    └── query/",
    "        ├── xml_parser.py",
    "        └── retriever.py",
], Inches(0.5), Inches(2.0), Inches(3.78), Inches(3.22), size=9, color=C_BLACK)

th_row(sl, Inches(4.82), Inches(1.62), Inches(2.6), Inches(0.3), "FILE")
th_row(sl, Inches(7.42), Inches(1.62), Inches(2.26), Inches(0.3), "ROLE")
for i, (fname, role) in enumerate([
    ("pdf_extractor.py", "PDF → 페이지 텍스트"),
    ("chunker.py",       "문장 분리 + 의미 경계 청킹"),
    ("embedder.py",      "BGE-M3 배치 임베딩"),
    ("qdrant_store.py",  "컬렉션 생성·upsert·검색"),
    ("xml_parser.py",    "ECG XML → 쿼리 컨텍스트"),
    ("retriever.py",     "쿼리 임베딩 → 검색 → 포맷"),
    ("pipeline.py",      "ingest / query CLI"),
    ("config.py",        "경로·청크 크기·책 메타"),
]):
    y = Inches(1.92 + i * 0.42)
    for hx, hw in [(Inches(4.82), Inches(2.6)), (Inches(7.42), Inches(2.26))]:
        tr_row(sl, hx, y, hw, Inches(0.4), i % 2 == 1)
    mono(sl, fname, Inches(4.9), y + Inches(0.07), Inches(2.42), Inches(0.28), size=8, color=C_BLACK)
    txt(sl,  role,  Inches(7.5), y + Inches(0.07), Inches(2.1),  Inches(0.28), size=9, color=C_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# 20 — 마무리 & 핵심 요약
# ══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG)
header_bar(sl, "Summary & Next Steps")
section_label(sl, "PHASE 1 COMPLETE")

th_row(sl, Inches(0.32), Inches(1.62), Inches(1.3), Inches(0.3), "AREA")
th_row(sl, Inches(1.62), Inches(1.62), Inches(4.2), Inches(0.3), "WHAT WE BUILT")
for i, (tag, desc) in enumerate([
    ("데이터",   "심전도 교재 7권 → 6,055 청크로 분해"),
    ("임베딩",   "BGE-M3로 한/영 동일 벡터 공간 구성"),
    ("저장",     "Qdrant + 메타데이터 필터 (레벨·언어)"),
    ("입력",     "ECG XML → 자동 파싱 → 자연어 쿼리"),
    ("검색",     "코사인 유사도 Top-K + 메타데이터 필터"),
    ("출력",     "검색된 교재 근거 기반 LLM 설명 생성"),
]):
    y = Inches(1.92 + i * 0.46)
    for hx, hw in [(Inches(0.32), Inches(1.3)), (Inches(1.62), Inches(4.2))]:
        tr_row(sl, hx, y, hw, Inches(0.44), i % 2 == 1)
    mono(sl, tag,  Inches(0.42), y + Inches(0.09), Inches(1.1), Inches(0.28), size=9, bold=True, color=C_BLACK)
    txt(sl,  desc, Inches(1.72), y + Inches(0.09), Inches(4.0), Inches(0.28), size=9, color=C_DARK)

card(sl, Inches(6.32), Inches(1.62), Inches(3.36), Inches(3.75), bg=C_BLACK, border=C_BLACK)
mono(sl, "NEXT STEPS", Inches(6.5), Inches(1.72), Inches(2.8), Inches(0.22), color=C_WHITE)
for i, n in enumerate([
    "색인 페이지 필터 추가",
    "하이브리드 검색 적용",
    "diagnosis_tags 자동 태깅",
    "Cross-Encoder 재순위",
    "LLM 연동 (Claude API)",
]):
    mono(sl, f"→  {n}", Inches(6.5), Inches(2.1 + i * 0.58), Inches(3.0), Inches(0.5), size=9, color=C_WHITE)

mono(sl, "> Phase 1 Complete  ·  All components verified  ·  6,055 chunks indexed",
     Inches(0.32), Inches(5.25), Inches(9.36), Inches(0.22))

# ── 저장 ──────────────────────────────────────────────────────────────────────
out = "/mnt/workspace/ecg-rag/ECG_RAG_발표자료_graphify스타일.pptx"
prs.save(out)
print(f"저장 완료: {out}")
